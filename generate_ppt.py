from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colour palette ──────────────────────────────────────────────
BG_DARK      = RGBColor(0x06, 0x0a, 0x1a)   # deep navy
BG_CARD      = RGBColor(0x0d, 0x15, 0x2e)   # card navy
CYAN         = RGBColor(0x00, 0xe5, 0xff)   # CybeWatch cyan
PURPLE       = RGBColor(0x7c, 0x3a, 0xed)   # accent purple
WHITE        = RGBColor(0xff, 0xff, 0xff)
GREY         = RGBColor(0xaa, 0xbb, 0xcc)
GREEN        = RGBColor(0x00, 0xe5, 0x96)
ORANGE       = RGBColor(0xff, 0x85, 0x00)
RED_BADGE    = RGBColor(0xff, 0x3b, 0x3b)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

IMG_DASHBOARD = "assets/screenshots/dashboard-preview.png"
IMG_TEAM      = "assets/screenshots/team-page.png"

# ── helpers ─────────────────────────────────────────────────────
def bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h,
        size=20, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def badge(slide, label, l, t, color=CYAN, text_color=BG_DARK, size=11):
    bw = Inches(len(label)*0.085 + 0.3)
    bh = Inches(0.28)
    box(slide, l, t, bw, bh, color)
    txt(slide, label, l + Inches(0.08), t + Inches(0.03),
        bw, bh, size=size, bold=True, color=text_color, align=PP_ALIGN.LEFT)

def divider(slide, t, color=CYAN, alpha=80):
    shape = slide.shapes.add_shape(1, Inches(0.5), t, Inches(12.33), Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def header(slide, title, subtitle=None):
    """Top accent bar + title."""
    box(slide, 0, 0, W, Inches(0.08), CYAN)
    txt(slide, title,
        Inches(0.55), Inches(0.18), Inches(10), Inches(0.7),
        size=34, bold=True, color=CYAN)
    if subtitle:
        txt(slide, subtitle,
            Inches(0.55), Inches(0.82), Inches(10), Inches(0.45),
            size=15, color=GREY, italic=True)

def bullet_block(slide, title, points, l, t, w, h,
                 title_color=CYAN, bullet_color=WHITE,
                 badge_label=None, badge_color=CYAN):
    box(slide, l, t, w, h, BG_CARD)
    yt = t + Inches(0.15)
    if badge_label:
        badge(slide, badge_label, l + Inches(0.2), yt, badge_color)
        yt += Inches(0.35)
    txt(slide, title, l + Inches(0.2), yt, w - Inches(0.4), Inches(0.45),
        size=16, bold=True, color=title_color)
    yt += Inches(0.42)
    for pt in points:
        txt(slide, f"  •  {pt}", l + Inches(0.1), yt,
            w - Inches(0.3), Inches(0.38), size=12, color=bullet_color)
        yt += Inches(0.3)

# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0, 0, W, Inches(0.08), CYAN)
box(s, 0, H - Inches(0.08), W, Inches(0.08), PURPLE)

# Glow box behind title
box(s, Inches(1.5), Inches(1.6), Inches(10.3), Inches(2.8), RGBColor(0x0d,0x18,0x3a))

txt(s, "🛡️  CybeWatch Enterprise",
    Inches(2), Inches(1.8), Inches(9.5), Inches(1.2),
    size=46, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

txt(s, "AI-Driven Cybersecurity Platform — Real-Time Threat Detection & Response",
    Inches(2), Inches(2.9), Inches(9.5), Inches(0.6),
    size=18, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

divider(s, Inches(3.65))

txt(s, "Team Glanzee  |  Abhinav R.  &  SF Melena",
    Inches(2), Inches(3.8), Inches(9.5), Inches(0.5),
    size=16, color=GREY, align=PP_ALIGN.CENTER)

badge(s, "HACKATHON SUBMISSION", Inches(5.3), Inches(4.5),
      color=PURPLE, text_color=WHITE, size=13)

txt(s, "GitHub: github.com/abhinavrenjithg-sys/CYBEWATCH",
    Inches(2), Inches(5.5), Inches(9.5), Inches(0.4),
    size=13, color=CYAN, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "The Problem", "Why CybeWatch Exists")

problems = [
    ("⚠️  Reactive Security",   "Most orgs detect breaches AFTER damage is done — avg. 207 days to identify a breach."),
    ("💸  Cost Barrier",         "Enterprise SIEMs (Splunk, CrowdStrike) cost $50K–$500K/yr — out of reach for SMEs."),
    ("🔊  Alert Fatigue",        "Security teams drown in thousands of false-positive alerts daily, missing real threats."),
    ("🧩  Fragmented Tooling",   "No single platform combines scanning, network monitoring, log analysis, and AI response."),
]

y = Inches(1.5)
for icon_title, desc in problems:
    card = box(s, Inches(0.5), y, Inches(12.3), Inches(0.85), BG_CARD)
    txt(s, icon_title, Inches(0.75), y + Inches(0.1), Inches(3.2), Inches(0.6),
        size=15, bold=True, color=RED_BADGE)
    txt(s, desc, Inches(3.8), y + Inches(0.22), Inches(8.8), Inches(0.5),
        size=13, color=WHITE)
    y += Inches(1.0)

txt(s, "CybeWatch solves all four — in one unified, affordable, AI-powered platform.",
    Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.5),
    size=15, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — Solution Overview
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Our Solution — CybeWatch", "A fully integrated, real-time threat detection & response platform")

cols = [
    ("🔍", "Web Scanner",       ["Real HTTP/HTTPS scanning", "SSL certificate audit", "Header security checks", "Vulnerability reporting"], GREEN, "LIVE ✅"),
    ("📡", "Network Monitor",   ["Live packet interception", "ARP poisoning detection", "Port scan detection", "Traffic anomaly alerts"], CYAN,  "LIVE ✅"),
    ("📋", "Log Analyzer",      ["Syslog ingestion", "Event pattern correlation", "Anomaly flagging", "Timeline reconstruction"], PURPLE,"LIVE ✅"),
    ("🤖", "ML Threat Engine",  ["Random Forest classifier", "Gradient Boost pipeline", "DNN ensemble model", "99.98% accuracy (prototype)"], ORANGE, "PROTOTYPE ⚗️"),
]

x = Inches(0.3)
for icon, title, pts, color, badge_lbl in cols:
    card_w = Inches(3.1)
    card_h = Inches(4.8)
    box(s, x, Inches(1.5), card_w, card_h, BG_CARD)
    txt(s, icon, x + Inches(1.1), Inches(1.7), Inches(1), Inches(0.6), size=28, align=PP_ALIGN.CENTER)
    txt(s, title, x + Inches(0.1), Inches(2.35), card_w - Inches(0.2), Inches(0.5),
        size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
    bcol = GREEN if "LIVE" in badge_lbl else ORANGE
    badge(s, badge_lbl, x + Inches(0.6), Inches(2.82), color=bcol, text_color=BG_DARK, size=10)
    yp = Inches(3.2)
    for pt in pts:
        txt(s, f"• {pt}", x + Inches(0.15), yp, card_w - Inches(0.3), Inches(0.35),
            size=11, color=GREY)
        yp += Inches(0.32)
    x += Inches(3.28)

# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — Live Dashboard Screenshot
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Live Dashboard — Security Operations Center", "Real-time threat monitoring | LIVE ✅")

if os.path.exists(IMG_DASHBOARD):
    s.shapes.add_picture(IMG_DASHBOARD,
                         Inches(0.4), Inches(1.3),
                         Inches(12.5), Inches(5.7))
else:
    box(s, Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.7), BG_CARD)
    txt(s, "[Dashboard Screenshot]", Inches(5), Inches(3.5), Inches(3), Inches(1),
        size=18, color=GREY, align=PP_ALIGN.CENTER)

badge(s, "LIVE  ✅  Fully Functional", Inches(0.5), Inches(7.1),
      color=GREEN, text_color=BG_DARK, size=11)

# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — How It Works (Data Flow)
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "How It Works", "End-to-end data flow from raw events to actionable intelligence")

steps = [
    ("1", "COLLECT",   "Web Scanner\nNetwork Monitor\nLog Analyzer",          CYAN),
    ("2", "ENRICH",    "GeoIP tagging\nCVE lookups\nIOC correlation",          PURPLE),
    ("3", "DETECT",    "ML Ensemble\nRule-based SIEM\nAnomaly scoring",        ORANGE),
    ("4", "RESPOND",   "AI Playbooks\nARIA SOC Chat\nAuto-remediation",        GREEN),
    ("5", "PERSIST",   "Supabase DB\nReal-time sync\nVercel hosting",          CYAN),
]

x = Inches(0.3)
for num, title, body, color in steps:
    bw = Inches(2.45)
    box(s, x, Inches(1.6), bw, Inches(4.5), BG_CARD)
    # circle number
    circle = s.shapes.add_shape(9, x + Inches(0.83), Inches(1.85), Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    txt(s, num, x + Inches(0.88), Inches(1.92), Inches(0.7), Inches(0.6),
        size=18, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    txt(s, title, x + Inches(0.1), Inches(2.8), bw - Inches(0.2), Inches(0.45),
        size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s, body,  x + Inches(0.1), Inches(3.35), bw - Inches(0.2), Inches(1.5),
        size=11, color=GREY, align=PP_ALIGN.CENTER)
    # arrow (except last)
    if num != "5":
        txt(s, "→", x + bw + Inches(0.0), Inches(3.5), Inches(0.35), Inches(0.5),
            size=22, bold=True, color=GREY, align=PP_ALIGN.CENTER)
    x += bw + Inches(0.36)

# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — What's LIVE vs PROTOTYPE
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "What's Live vs Prototype", "Honest breakdown of current implementation status")

# LIVE column
box(s, Inches(0.4), Inches(1.4), Inches(5.9), Inches(5.7), BG_CARD)
badge(s, "✅  LIVE & FUNCTIONAL", Inches(0.7), Inches(1.55), color=GREEN, text_color=BG_DARK, size=12)
live_items = [
    "Real-time Dashboard with live metric cards",
    "Web Scanner — actual HTTP/HTTPS requests",
    "Network Monitor — live traffic analysis",
    "Log Analyzer — real syslog ingestion",
    "ARIA AI Chat — Gemini 2.0 Flash API",
    "Alerts engine with severity classification",
    "Investigation & Case Management UI",
    "Automated Playbook Builder",
    "Reports Center with export",
    "Supabase real-time database sync",
    "Kill Chain Topology Visualizer",
    "Global Attack Map (live GeoIP data)",
]
y = Inches(2.0)
for item in live_items:
    txt(s, f"✓  {item}", Inches(0.65), y, Inches(5.4), Inches(0.35),
        size=11.5, color=WHITE)
    y += Inches(0.34)

# PROTOTYPE column
box(s, Inches(6.8), Inches(1.4), Inches(6.1), Inches(5.7), BG_CARD)
badge(s, "⚗️  PROTOTYPE / DEMO", Inches(7.1), Inches(1.55), color=ORANGE, text_color=BG_DARK, size=12)
proto_items = [
    ("ML Threat Engine (train_model.py)",
     "Architecture complete. Model trained on KDD-Cup dataset. In production, would ingest live telemetry streams."),
    ("Neuro-Ensemble Classifier",
     "Random Forest + Gradient Boost + DNN pipeline built & tested. Accuracy: 99.98% on test set (simulated data)."),
    ("Dark Web OSINT Scanner",
     "UI and query interface built. Live Tor/onion scanning requires deployment on dedicated infrastructure."),
    ("Red Team Breach Simulator",
     "Simulation logic functional. Real exploit execution disabled for hackathon environment safety."),
    ("3D Global Attack Globe",
     "Three.js renderer working. Attack data is enriched with real GeoIP but animated for demo purposes."),
]
y = Inches(2.05)
for title, desc in proto_items:
    txt(s, f"⚗  {title}", Inches(7.05), y, Inches(5.6), Inches(0.35),
        size=12, bold=True, color=ORANGE)
    y += Inches(0.3)
    txt(s, desc, Inches(7.15), y, Inches(5.5), Inches(0.55),
        size=10, color=GREY, italic=True)
    y += Inches(0.65)

# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — AI SOC Assistant — ARIA
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "ARIA — AI SOC Assistant", "Powered by Google Gemini 2.0 Flash  |  LIVE ✅")

box(s, Inches(0.4), Inches(1.45), Inches(12.5), Inches(5.65), BG_CARD)

txt(s, "🤖", Inches(0.8), Inches(1.7), Inches(1), Inches(1), size=40)

txt(s, "What ARIA Does:",
    Inches(2), Inches(1.7), Inches(10), Inches(0.5),
    size=18, bold=True, color=CYAN)

aria_points = [
    "Accepts natural language security questions from the SOC analyst",
    "Analyses live threat data from the dashboard and provides context-aware responses",
    "Generates step-by-step remediation playbooks for detected threats",
    "Explains CVEs, attack techniques (MITRE ATT&CK), and recommended countermeasures",
    "Summarises daily/weekly threat reports on demand",
    "Integrated directly into the dashboard — no external tool switch needed",
]
y = Inches(2.35)
for pt in aria_points:
    txt(s, f"  ●  {pt}", Inches(2), y, Inches(10.5), Inches(0.38),
        size=13, color=WHITE)
    y += Inches(0.38)

badge(s, "LIVE ✅  Gemini 2.0 Flash API", Inches(0.7), Inches(6.6),
      color=GREEN, text_color=BG_DARK, size=11)
txt(s, "Note: API key required at runtime. Demo uses a sandboxed Gemini endpoint.",
    Inches(3.5), Inches(6.65), Inches(9), Inches(0.4),
    size=11, color=GREY, italic=True)

# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — Tech Stack
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Technology Stack", "Built with production-grade tools")

categories = [
    ("⚡  Runtime",        "Electron + Node.js",                   CYAN),
    ("🎨  Frontend",       "HTML5, Vanilla CSS3, ES6+ JS",         PURPLE),
    ("📊  Visualisation",  "Chart.js 4.4+, Three.js (3D Globe)",   CYAN),
    ("🤖  AI / ML",        "Gemini 2.0 Flash, Scikit-Learn, Pandas, NumPy", ORANGE),
    ("🗄️  Database",       "Supabase (PostgreSQL + Realtime)",     GREEN),
    ("☁️  Hosting",        "Vercel Serverless + Edge Functions",   CYAN),
    ("🌐  Network Intel",  "ip-api.com (GeoIP), CVE NVD API",     PURPLE),
    ("🔐  Auth",           "Supabase Auth — JWT + RLS Policies",   GREEN),
]

x = Inches(0.4)
y = Inches(1.5)
col = 0
for label, value, color in categories:
    box(s, x, y, Inches(6.1), Inches(0.78), BG_CARD)
    txt(s, label, x + Inches(0.2), y + Inches(0.05),
        Inches(2.5), Inches(0.65), size=12, bold=True, color=color)
    txt(s, value,  x + Inches(2.5), y + Inches(0.22),
        Inches(3.4), Inches(0.45), size=12, color=WHITE)
    col += 1
    if col % 2 == 0:
        x = Inches(0.4)
        y += Inches(0.92)
    else:
        x = Inches(6.8)

# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — Team
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Team Glanzee", "Founded 2024 — Two vibe coders on a mission to redefine enterprise security")

if os.path.exists(IMG_TEAM):
    s.shapes.add_picture(IMG_TEAM,
                         Inches(0.4), Inches(1.3),
                         Inches(12.5), Inches(5.7))
else:
    for xi, (name, role, skills, link) in enumerate([
        ("Abhinav R.", "Founder & Lead Architect", "Vibe Coder | Creative Thinker | UI/UX",
         "github.com/abhinavrenjithg-sys"),
        ("SF Melena",  "Co-Founder & UI/UX Director", "Vibe Coder | Creative Thinker | Full-Stack",
         "Team Glanzee"),
    ]):
        cx = Inches(1.0) + xi * Inches(6.5)
        box(s, cx, Inches(1.6), Inches(5.5), Inches(5.0), BG_CARD)
        txt(s, name,  cx + Inches(0.3), Inches(2.0),
            Inches(5), Inches(0.6), size=22, bold=True, color=WHITE)
        txt(s, role,  cx + Inches(0.3), Inches(2.65),
            Inches(5), Inches(0.45), size=13, color=CYAN)
        txt(s, skills, cx + Inches(0.3), Inches(3.2),
            Inches(5), Inches(0.4), size=12, color=GREY, italic=True)
        txt(s, link,  cx + Inches(0.3), Inches(3.75),
            Inches(5), Inches(0.35), size=11, color=PURPLE)

# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — Conclusion & GitHub
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0, 0, W, Inches(0.08), CYAN)
box(s, 0, H - Inches(0.08), W, Inches(0.08), PURPLE)

txt(s, "Conclusion",
    Inches(1), Inches(0.4), Inches(11), Inches(0.8),
    size=36, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

conclusion = (
    "CybeWatch demonstrates that enterprise-grade cybersecurity is achievable without the cost and complexity "
    "of traditional SIEM platforms. By combining real-time system monitoring, a Neuro-Ensemble ML threat engine, "
    "and an AI-powered SOC assistant into a single unified platform, CybeWatch delivers actionable security "
    "intelligence at the speed modern organizations demand.\n\n"
    "The platform is fully functional, cloud-synced, and production-ready — with clearly marked prototype "
    "components that outline the path to full deployment. Team Glanzee built CybeWatch with the conviction "
    "that every organization deserves the tools to defend itself."
)
box(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(3.2), BG_CARD)
txt(s, conclusion,
    Inches(0.9), Inches(1.5), Inches(11.5), Inches(3.0),
    size=14, color=WHITE, align=PP_ALIGN.LEFT)

divider(s, Inches(4.75))

txt(s, "🔗  GitHub Repository",
    Inches(1), Inches(4.95), Inches(11), Inches(0.5),
    size=16, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

txt(s, "https://github.com/abhinavrenjithg-sys/CYBEWATCH",
    Inches(1), Inches(5.5), Inches(11), Inches(0.5),
    size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "Public repository  •  Open source  •  All commits timestamped",
    Inches(1), Inches(6.1), Inches(11), Inches(0.4),
    size=12, color=GREY, align=PP_ALIGN.CENTER, italic=True)

# ── Save ──────────────────────────────────────────────────────────
OUTPUT = "CybeWatch_Presentation.pptx"
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}  ({prs.slides.__len__()} slides)")
